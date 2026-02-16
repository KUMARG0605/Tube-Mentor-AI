"""Slides Generator Service - Create PowerPoint presentations using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import List, Optional
import re

from app.config import settings

# Output directory for slides
SLIDES_DIR = settings.OUTPUT_DIR / "slides"
SLIDES_DIR.mkdir(parents=True, exist_ok=True)


# Color scheme
COLORS = {
    "primary": RgbColor(139, 92, 246),      # Purple (brand color)
    "secondary": RgbColor(30, 41, 59),       # Dark slate
    "accent": RgbColor(236, 72, 153),        # Pink
    "background": RgbColor(15, 23, 42),      # Dark background
    "text": RgbColor(248, 250, 252),         # Light text
    "muted": RgbColor(148, 163, 184),        # Muted gray
}


def create_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["background"]
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS["muted"]
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_content_slide(
    prs: Presentation, 
    title: str, 
    bullet_points: List[str],
    slide_number: int = None
):
    """Create a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["background"]
    background.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["secondary"]
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    
    # Bullet points
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.6), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["text"]
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    # Slide number
    if slide_number:
        num_box = slide.shapes.add_textbox(
            Inches(9), Inches(6.8), Inches(0.8), Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_number)
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["muted"]
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def create_section_slide(prs: Presentation, section_title: str, section_number: int):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["secondary"]
    background.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"SECTION {section_number}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_conclusion_slide(prs: Presentation, key_takeaways: List[str]):
    """Create a conclusion/takeaways slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["background"]
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # Takeaways
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.8), Inches(8.6), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(key_takeaways):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS["text"]
        p.space_before = Pt(16)
    
    return slide


def parse_summary_to_slides(summary_text: str) -> dict:
    """Parse a summary into slide-friendly sections."""
    sections = []
    current_section = {"title": "", "points": []}
    
    lines = summary_text.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for headers (## or ** wrapped)
        if line.startswith('## ') or line.startswith('# '):
            if current_section["title"] and current_section["points"]:
                sections.append(current_section)
            current_section = {"title": line.lstrip('#').strip(), "points": []}
        elif line.startswith('**') and line.endswith('**'):
            if current_section["title"] and current_section["points"]:
                sections.append(current_section)
            current_section = {"title": line.strip('*').strip(), "points": []}
        elif line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
            point = line.lstrip('-*• ').strip()
            if point and len(point) > 3:
                current_section["points"].append(point)
        elif line.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
            point = re.sub(r'^\d+\.\s*', '', line).strip()
            if point and len(point) > 3:
                current_section["points"].append(point)
    
    # Add last section
    if current_section["title"] and current_section["points"]:
        sections.append(current_section)
    
    return sections


def generate_presentation(
    title: str,
    summary: str,
    key_topics: List[str] = None,
    video_id: str = None
) -> dict:
    """Generate a PowerPoint presentation from a video summary."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse summary into sections
    sections = parse_summary_to_slides(summary)
    
    # Title slide
    create_title_slide(prs, title, "Video Summary Presentation")
    
    # Table of contents if we have topics
    if key_topics and len(key_topics) > 0:
        create_content_slide(prs, "Topics Covered", key_topics[:8], slide_number=2)
    
    # Content slides
    slide_num = 3
    for i, section in enumerate(sections):
        if len(section["points"]) > 0:
            # Split into multiple slides if too many points
            points = section["points"]
            chunks = [points[i:i+5] for i in range(0, len(points), 5)]
            
            for j, chunk in enumerate(chunks):
                section_title = section["title"]
                if len(chunks) > 1:
                    section_title = f"{section['title']} ({j+1}/{len(chunks)})"
                create_content_slide(prs, section_title, chunk, slide_number=slide_num)
                slide_num += 1
    
    # Conclusion slide with key takeaways
    takeaways = key_topics[:5] if key_topics else ["Review the content", "Practice what you learned"]
    create_conclusion_slide(prs, takeaways)
    
    # Save presentation
    filename = f"presentation_{video_id or 'custom'}.pptx"
    filepath = SLIDES_DIR / filename
    prs.save(str(filepath))
    
    return {
        "filename": filename,
        "filepath": str(filepath),
        "slide_count": len(prs.slides),
        "title": title,
    }


def generate_custom_presentation(
    title: str,
    slides_content: List[dict]
) -> dict:
    """Generate a custom presentation from provided content.
    
    slides_content: List of {"title": str, "points": List[str]}
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    create_title_slide(prs, title, "Created with TubeMentor AI")
    
    # Content slides
    for i, slide_data in enumerate(slides_content):
        create_content_slide(
            prs, 
            slide_data.get("title", f"Slide {i+1}"),
            slide_data.get("points", []),
            slide_number=i+2
        )
    
    # Save
    import uuid
    filename = f"custom_{uuid.uuid4().hex[:8]}.pptx"
    filepath = SLIDES_DIR / filename
    prs.save(str(filepath))
    
    return {
        "filename": filename,
        "filepath": str(filepath),
        "slide_count": len(prs.slides),
        "title": title,
    }
